#!/usr/bin/env python3
"""pdf2pptx — PDFを編集可能なPPTXに変換する (Phase 1.1).

方式:
  1. PyMuPDF でテキストを span/line/block 単位で抽出する（横書き・可視のみ編集対象）
  2. 編集対象テキストの領域を redaction で背景から削除する（二重写り防止。
     画像・罫線・図形は残す）
  3. redaction 後のページを画像化してスライド背景に敷く
  4. 抽出テキストを編集可能なテキストボックスとして同じ座標に重ねる

Phase 1 では安全側に倒し、次のケースは編集対象にせず背景画像側に残す:
  - 縦書き・回転テキスト（wmode / dir ベクトルで判定）
  - ページ自体が回転している場合（/Rotate。座標補正はPhase 2以降）
  - 不可視テキスト（alpha=0 または render mode 3。OCRの不可視テキスト層等）
  - 文字化けテキスト（ToUnicode欠落等）
  - 上記いずれかと bbox が重なる横書きテキスト（redactionによる欠損防止）
  - ページ面積の大部分（既定85%以上）が画像で覆われている場合、そのページの
    可視テキストも含めて編集対象にしない（画像に焼き込まれた文字を編集可能
    テキスト化すると、画像と重なって二重写りになるリスクがあるため）

Web公開等を見据え、ページ数・dpi・画素数（1ページ/全ページ合計）・入出力ファイル
サイズ・処理時間に上限を設けている。超過時は変換前、または各ページ処理開始時点で
安全に停止する（ソフトタイムアウト）。プロセスが応答しなくなった場合に確実に
強制終了したい場合は worker.py（別プロセス+ハードタイムアウト）を使うこと。

使い方:
  python convert.py input.pdf output.pptx [--dpi 150] [--debug-dir DIR]
"""

from __future__ import annotations

import argparse
import csv
import io
import os
import re
import shutil
import subprocess
import sys
import tempfile
import time
from dataclasses import dataclass, field
from pathlib import Path

import pymupdf
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

EMU_PER_PT = 12700
PT_PER_INCH = 72.0

# 横書き判定: line の書字方向ベクトルが (1, 0) にほぼ一致すること
HORIZONTAL_DIR_TOL = 0.02
# redaction 矩形を少し縮めて隣接コンテンツの巻き込みを防ぐ (pt)
REDACT_SHRINK = 0.2
# span 中の置換文字 (U+FFFD) がこの割合を超えたら文字化けとみなし背景に残す
GARBLED_RATIO = 0.3
# テキストボックスの幅の余裕。代替フォントの字幅差で右端が欠けるのを防ぐ (pt)
BOX_PAD = 1.0
# ページサイズ差をこの値 (pt) まで「同一サイズ」とみなす
PAGE_SIZE_TOL = 1.0
# ページ面積のこの割合以上を画像が覆っていたら「全面画像ページ」とみなし、
# 可視テキストがあっても編集対象にしない（二重写りリスクの回避）
FULL_IMAGE_COVERAGE_THRESHOLD = 0.85
# PowerPoint (OOXML) のスライドサイズ制限
PPTX_MIN_SLIDE_INCH = 1.0
PPTX_MAX_SLIDE_INCH = 56.0

# --- Web公開等を見据えた基本制限（既定値。CLI/API引数で上書き可能） ---
DEFAULT_MAX_PAGES = 200
DEFAULT_MAX_DPI = 300
DEFAULT_MAX_PAGE_PIXELS = 50_000_000     # 背景画像1ページあたりの画素数上限（約50MP）
DEFAULT_MAX_TOTAL_PIXELS = 500_000_000   # 背景画像の全ページ合計の画素数上限（約500MP）
DEFAULT_MAX_FILE_SIZE_MB = 100.0
DEFAULT_MAX_OUTPUT_SIZE_MB = 300.0

# --- OCR（任意機能・既定OFF） ---
# スキャンPDFを編集可能テキスト化するにはOCRが必要だが、誤認識・二重写り・
# 依存環境差が大きいため、既定では無効。ローカルPCにTesseractが入っている
# 場合だけ、明示オプションで有効化する。
OCR_ENGINE_OFF = "off"
OCR_ENGINE_TESSERACT = "tesseract"
OCR_ENGINES = (OCR_ENGINE_OFF, OCR_ENGINE_TESSERACT)
DEFAULT_OCR_ENGINE = OCR_ENGINE_OFF
DEFAULT_OCR_LANG = "jpn+eng"
DEFAULT_OCR_MIN_CONF = 35.0
DEFAULT_OCR_TIMEOUT = 30.0


# ---------------------------------------------------------------------------
# フォントマッピング
# ---------------------------------------------------------------------------

# PDF内フォント名（サブセット接頭辞除去・小文字・空白/ハイフン除去後）の
# 完全一致マップ。ヒューリスティックより優先される。
FONT_MAP = {
    "msgothic": "ＭＳ ゴシック",
    "mspgothic": "ＭＳ Ｐゴシック",
    "msmincho": "ＭＳ 明朝",
    "mspmincho": "ＭＳ Ｐ明朝",
    "meiryo": "Meiryo",
    "meiryoui": "Meiryo UI",
    "yugothic": "Yu Gothic",
    "yumincho": "Yu Mincho",
    "arial": "Arial",
    "helvetica": "Arial",
    "timesnewroman": "Times New Roman",
    "times": "Times New Roman",
    "couriernew": "Courier New",
    "courier": "Courier New",
    "calibri": "Calibri",
}

JP_HINTS = (
    "gothic", "mincho", "meiryo", "hiragino", "hira", "yugo", "yumin",
    "kozuka", "kozgo", "kozmin", "sourcehan", "notosanscjk", "notosansjp",
    "notoserifcjk", "notoserifjp", "ipaex", "ipag", "ipam", "biz", "udshingo",
    "ryumin", "shingo", "midashi", "maru", "kaku", "japan",
)
SERIF_HINTS = (
    "mincho", "ryumin", "kozmin", "notoserif", "sourcehanserif", "ipam",
    "times", "georgia", "garamond", "serif", "roman", "century", "cambria",
    "palatino", "minion", "book",
)
MONO_HINTS = ("courier", "mono", "consolas", "menlo")

FALLBACK_JP_SANS = "Yu Gothic"
FALLBACK_JP_SERIF = "Yu Mincho"
FALLBACK_LATIN_SANS = "Arial"
FALLBACK_LATIN_SERIF = "Times New Roman"
FALLBACK_MONO = "Courier New"

# span flags (PyMuPDF: フォントの見た目上のスタイル)
FLAG_ITALIC = 1 << 1
FLAG_MONO = 1 << 3
FLAG_BOLD = 1 << 4

# span char_flags (PyMuPDF/MuPDF: FZ_STEXT_*。描画のされ方を表す別のビット集合)
CHAR_FLAG_FILLED = 16   # FZ_STEXT_FILLED
CHAR_FLAG_STROKED = 32  # FZ_STEXT_STROKED

_SUBSET_PREFIX = re.compile(r"^[A-Z]{6}\+")


def _contains_cjk(text: str) -> bool:
    return any(
        "　" <= ch <= "ヿ"      # 記号・かな
        or "一" <= ch <= "鿿"   # 漢字
        or "＀" <= ch <= "￯"   # 全角英数・半角カナ
        for ch in text
    )


_LATIN_ONLY_FAMILIES = {"Arial", "Times New Roman", "Courier New", "Calibri"}


def map_font(pdf_font_name: str, text: str) -> str:
    """PDF内のフォント名を PowerPoint で使えるファミリー名に変換する。"""
    name = _SUBSET_PREFIX.sub("", pdf_font_name or "")
    key = name.lower().replace(" ", "").replace("-", "").replace(",", "")
    for exact, family in FONT_MAP.items():
        if key.startswith(exact):
            # 日本語テキストが欧文専用フォント名 (Arial Unicode MS 等) に
            # 載っている場合は日本語フォールバックを優先する
            if family in _LATIN_ONLY_FAMILIES and _contains_cjk(text):
                break
            return family
    is_jp = _contains_cjk(text) or any(h in key for h in JP_HINTS)
    is_serif = any(h in key for h in SERIF_HINTS)
    if is_jp:
        return FALLBACK_JP_SERIF if is_serif else FALLBACK_JP_SANS
    if any(h in key for h in MONO_HINTS):
        return FALLBACK_MONO
    return FALLBACK_LATIN_SERIF if is_serif else FALLBACK_LATIN_SANS


def _is_bold(span: dict) -> bool:
    if span["flags"] & FLAG_BOLD:
        return True
    name = (span.get("font") or "").lower()
    if any(h in name for h in ("bold", "heavy", "black", "semibold", "demibold")):
        return True
    # ヒラギノ等のウェイト表記 (W6以上を太字扱い)
    m = re.search(r"w(\d)\b", name)
    return bool(m and int(m.group(1)) >= 6)


def _is_italic(span: dict) -> bool:
    if span["flags"] & FLAG_ITALIC:
        return True
    name = (span.get("font") or "").lower()
    return "italic" in name or "oblique" in name


def _is_invisible(span: dict) -> bool:
    """描画されない(見えない)テキストか判定する。OCRの不可視テキスト層等。

    - alpha == 0: 完全に透明な塗り（ウォーターマーク・隠しテキスト等）
    - char_flags に FILLED も STROKED も立っていない: PDFのテキスト
      描画モード3（invisible）。OCRソフトが原稿画像の上に検索用テキストを
      重ねる際の定番の手法。
    """
    if span.get("alpha", 255) == 0:
        return True
    char_flags = span.get("char_flags", CHAR_FLAG_FILLED)
    return not (char_flags & (CHAR_FLAG_FILLED | CHAR_FLAG_STROKED))


# ---------------------------------------------------------------------------
# 抽出
# ---------------------------------------------------------------------------

@dataclass
class Line:
    """編集対象として抽出した1行分のテキスト。"""
    bbox: pymupdf.Rect
    spans: list[dict]


@dataclass
class PageResult:
    width: float                # pt
    height: float               # pt
    lines: list[Line] = field(default_factory=list)
    image_png: bytes = b""
    warnings: list[str] = field(default_factory=list)


def _is_garbled(text: str) -> bool:
    if not text:
        return False
    return text.count("�") / len(text) > GARBLED_RATIO


# ---------------------------------------------------------------------------
# OCR（任意機能）
# ---------------------------------------------------------------------------

def _normalize_ocr_engine(engine: str) -> str:
    engine = (engine or OCR_ENGINE_OFF).lower()
    if engine not in OCR_ENGINES:
        raise ValueError(
            f"ocr は {', '.join(OCR_ENGINES)} のいずれかである必要があります: {engine}"
        )
    return engine


def _tesseract_path() -> str | None:
    return shutil.which("tesseract")


def _split_ocr_langs(lang: str) -> list[str]:
    return [part for part in (lang or "").split("+") if part]


def check_ocr_available(engine: str, lang: str = DEFAULT_OCR_LANG) -> None:
    """OCRエンジンが利用可能かを軽く検査する。

    OCRは任意機能なので、offの場合は何もしない。Tesseractを使う場合は
    コマンド本体と指定言語データが見つかるかを事前確認し、変換途中で
    分かりにくく失敗しないようにする。
    """
    engine = _normalize_ocr_engine(engine)
    if engine == OCR_ENGINE_OFF:
        return

    exe = _tesseract_path()
    if not exe:
        raise ValueError(
            "OCRを使うにはTesseractのインストールが必要です。"
            "通常変換または高画質変換はOCRなしでも利用できます"
        )

    requested = set(_split_ocr_langs(lang))
    if not requested:
        raise ValueError("ocr_lang は空にできません")

    try:
        res = subprocess.run(
            [exe, "--list-langs"],
            capture_output=True,
            text=True,
            timeout=5,
            check=False,
        )
    except (OSError, subprocess.TimeoutExpired):
        # Tesseract本体の存在だけは確認済み。言語一覧取得に失敗しても、
        # 実OCR時のエラーに委ねる。
        return

    available = {
        line.strip()
        for line in res.stdout.splitlines()
        if line.strip() and not line.lower().startswith("list of available")
    }
    missing = sorted(requested - available)
    if missing:
        raise ValueError(
            "OCRに必要なTesseract言語データが見つかりません: "
            + ", ".join(missing)
        )


def _parse_conf(value: str) -> float:
    try:
        return float(value)
    except (TypeError, ValueError):
        return -1.0


def _join_ocr_parts(parts: list[str]) -> str:
    """OCR TSVの単語断片を行テキストへまとめる。

    英単語同士は空白でつなぎ、日本語などCJKを含む断片は空白を挟まずに
    つなぐ。Tesseractの日本語出力は分かち書きが不安定なため、安全側の
    簡易ヒューリスティックにしている。
    """
    out = ""
    for part in parts:
        part = part.strip()
        if not part:
            continue
        if not out:
            out = part
            continue
        if _contains_cjk(out[-1]) or _contains_cjk(part[0]):
            out += part
        else:
            out += " " + part
    return out


def _run_tesseract_tsv(
    image_png: bytes,
    *,
    lang: str,
    timeout: float,
) -> str:
    exe = _tesseract_path()
    if not exe:
        raise RuntimeError("Tesseractが見つかりません")

    with tempfile.TemporaryDirectory(prefix="pdf2pptx_ocr_") as tmp:
        image_path = Path(tmp) / "page.png"
        image_path.write_bytes(image_png)
        res = subprocess.run(
            [exe, str(image_path), "stdout", "-l", lang, "--psm", "6", "tsv"],
            capture_output=True,
            text=True,
            timeout=timeout,
            check=False,
        )
    if res.returncode != 0:
        detail = (res.stderr or res.stdout or "Tesseract OCR failed").strip()
        raise RuntimeError(f"Tesseract OCRに失敗しました: {detail}")
    return res.stdout


def _ocr_lines_from_tsv(
    tsv_text: str,
    *,
    dpi: int,
    min_conf: float,
) -> list[Line]:
    """Tesseract TSVを既存のLine形式へ変換する。

    Tesseractの座標は画像ピクセル基準。render_background()と同じdpiで
    作った画像をOCRするため、72/dpiでPDFポイント座標へ戻せる。
    """
    scale = PT_PER_INCH / dpi
    grouped: dict[tuple[str, str, str], list[tuple[pymupdf.Rect, str]]] = {}
    reader = csv.DictReader(io.StringIO(tsv_text), delimiter="\t")
    for row in reader:
        if row.get("level") != "5":
            continue
        text = (row.get("text") or "").strip()
        if not text:
            continue
        conf = _parse_conf(row.get("conf", ""))
        if conf >= 0 and conf < min_conf:
            continue
        try:
            left = float(row["left"])
            top = float(row["top"])
            width = float(row["width"])
            height = float(row["height"])
        except (KeyError, ValueError):
            continue
        if width <= 0 or height <= 0:
            continue
        key = (
            row.get("block_num", "0"),
            row.get("par_num", "0"),
            row.get("line_num", "0"),
        )
        rect = pymupdf.Rect(
            left * scale,
            top * scale,
            (left + width) * scale,
            (top + height) * scale,
        )
        grouped.setdefault(key, []).append((rect, text))

    lines: list[Line] = []
    for items in grouped.values():
        if not items:
            continue
        bbox = pymupdf.Rect(items[0][0])
        parts: list[str] = []
        for rect, text in items:
            bbox |= rect
            parts.append(text)
        line_text = _join_ocr_parts(parts)
        if not line_text.strip():
            continue
        font_size = max(min(bbox.height * 0.78, bbox.height), 4.0)
        lines.append(
            Line(
                bbox=bbox,
                spans=[
                    {
                        "text": line_text,
                        "bbox": tuple(bbox),
                        "font": FALLBACK_JP_SANS,
                        "size": font_size,
                        "color": 0x000000,
                        "flags": 0,
                    }
                ],
            )
        )
    return lines


def ocr_image_to_lines(
    image_png: bytes,
    *,
    dpi: int,
    lang: str = DEFAULT_OCR_LANG,
    min_conf: float = DEFAULT_OCR_MIN_CONF,
    timeout: float = DEFAULT_OCR_TIMEOUT,
) -> list[Line]:
    tsv_text = _run_tesseract_tsv(image_png, lang=lang, timeout=timeout)
    return _ocr_lines_from_tsv(tsv_text, dpi=dpi, min_conf=min_conf)


# get_text("dict") の既定フラグ (TEXTFLAGS_DICT) は画像ブロックの埋め込み
# (TEXT_PRESERVE_IMAGES) を含む。テキストブロックしか使わないため外し、
# 画像が多いPDFで不要なバイナリをメモリに載せないようにする。
_DICT_FLAGS_NO_IMAGES = pymupdf.TEXTFLAGS_DICT & ~pymupdf.TEXT_PRESERVE_IMAGES


def _get_text_dict(page: pymupdf.Page) -> dict:
    return page.get_text("dict", flags=_DICT_FLAGS_NO_IMAGES)


def _text_visibility(text_dict: dict) -> tuple[bool, bool]:
    """(可視テキストが1文字でもあるか, 不可視テキストが1文字でもあるか) を返す。"""
    any_visible = False
    any_invisible = False
    for block in text_dict["blocks"]:
        if block["type"] != 0:
            continue
        for line in block["lines"]:
            for span in line["spans"]:
                if not span["text"].strip():
                    continue
                if _is_invisible(span):
                    any_invisible = True
                else:
                    any_visible = True
    return any_visible, any_invisible


def _image_coverage_ratio(page: pymupdf.Page) -> float:
    """ページ面積に対する画像の被覆率を概算する(0.0〜1.0)。

    厳密な和集合面積は計算せず、「最大の1枚の画像の面積」と「全画像面積の
    単純合計（重なりを考慮しない概算）」の大きい方を採用する。全面スキャン画像
    (1枚が全面を覆う)・タイル状に分割された全面画像のどちらでも検出できる、
    安全側に倒した近似。
    """
    page_area = page.rect.width * page.rect.height
    if page_area <= 0:
        return 0.0
    infos = page.get_image_info()
    if not infos:
        return 0.0
    largest = 0.0
    total = 0.0
    for info in infos:
        r = pymupdf.Rect(info["bbox"]) & page.rect
        area = max(r.width, 0.0) * max(r.height, 0.0)
        largest = max(largest, area)
        total += area
    return min(max(largest, total) / page_area, 1.0)


def extract_editable_lines(
    page: pymupdf.Page, text_dict: dict | None = None
) -> tuple[list[Line], list[str]]:
    """横書き・可視・重なりのない行だけを編集対象として抽出する。

    縦書き・回転・不可視(OCR等)・文字化けのテキストは背景画像に残す。
    それらと bbox が重なる横書き行も、redactionによる欠損を防ぐため
    背景画像に残す（該当行ごと編集対象から除外する）。
    """
    d = text_dict if text_dict is not None else _get_text_dict(page)

    candidates: list[Line] = []
    fallback_bboxes: list[pymupdf.Rect] = []
    n_vertical = n_invisible = n_garbled = 0

    for block in d["blocks"]:
        if block["type"] != 0:
            continue
        for raw_line in block["lines"]:
            is_horizontal = raw_line.get("wmode", 0) == 0
            if is_horizontal:
                dx, dy = raw_line["dir"]
                is_horizontal = (
                    dx > 1.0 - HORIZONTAL_DIR_TOL and abs(dy) < HORIZONTAL_DIR_TOL
                )
            if not is_horizontal:
                for span in raw_line["spans"]:
                    if span["text"].strip():
                        n_vertical += 1
                        fallback_bboxes.append(pymupdf.Rect(span["bbox"]))
                continue

            spans = []
            for span in raw_line["spans"]:
                if not span["text"].strip():
                    continue
                if _is_invisible(span):
                    n_invisible += 1
                    fallback_bboxes.append(pymupdf.Rect(span["bbox"]))
                    continue
                if _is_garbled(span["text"]):
                    n_garbled += 1
                    fallback_bboxes.append(pymupdf.Rect(span["bbox"]))
                    continue
                spans.append(span)
            if not spans:
                continue
            bbox = pymupdf.Rect(spans[0]["bbox"])
            for span in spans[1:]:
                bbox |= pymupdf.Rect(span["bbox"])
            candidates.append(Line(bbox=bbox, spans=spans))

    lines: list[Line] = []
    n_overlap = 0
    for line in candidates:
        if any(line.bbox.intersects(fb) for fb in fallback_bboxes):
            n_overlap += 1
            continue
        lines.append(line)

    warnings: list[str] = []
    if n_vertical:
        warnings.append(
            f"縦書き・回転テキスト {n_vertical} 行は編集対象外です（背景画像に残します）"
        )
    if n_invisible:
        warnings.append(
            f"不可視のテキスト(OCR等) {n_invisible} 箇所は編集対象外です（背景画像に残します）"
        )
    if n_garbled:
        warnings.append(
            f"文字コードを復元できないテキスト {n_garbled} 箇所を背景画像に残しました"
        )
    if n_overlap:
        warnings.append(
            f"他の要素と重なる横書きテキスト {n_overlap} 行を、欠損防止のため背景画像に残しました"
        )
    return lines, warnings


def redact_text(page: pymupdf.Page, lines: list[Line]) -> None:
    """編集対象テキストを背景から削除する。画像・罫線・図形は残す。"""
    s = REDACT_SHRINK
    for line in lines:
        for span in line.spans:
            r = pymupdf.Rect(span["bbox"])
            if r.width > 2 * s and r.height > 2 * s:
                r = pymupdf.Rect(r.x0 + s, r.y0 + s, r.x1 - s, r.y1 - s)
            # fill=False: 塗り潰さず、下にある画像・図形をそのまま見せる
            page.add_redact_annot(r, fill=False)
    if not lines:
        return
    try:
        page.apply_redactions(
            images=pymupdf.PDF_REDACT_IMAGE_NONE,
            graphics=pymupdf.PDF_REDACT_LINE_ART_NONE,
        )
    except TypeError:  # 古いPyMuPDF (graphics引数なし)
        page.apply_redactions(images=pymupdf.PDF_REDACT_IMAGE_NONE)


def render_background(page: pymupdf.Page, dpi: int) -> bytes:
    """redaction適用後のページをPNG化する。"""
    pix = page.get_pixmap(dpi=dpi, alpha=False)
    img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
    buf = io.BytesIO()
    img.save(buf, format="PNG", optimize=True)
    return buf.getvalue()


def _page_pixel_count(page: pymupdf.Page, dpi: int) -> float:
    """指定dpiで背景画像化した場合の概算画素数。"""
    w_px = page.rect.width / 72.0 * dpi
    h_px = page.rect.height / 72.0 * dpi
    return w_px * h_px


def _check_pixel_budget(
    page: pymupdf.Page, dpi: int, max_page_pixels: int, page_no: int
) -> None:
    """背景画像化した場合の画素数が1ページあたりの上限を超えないか検査する。"""
    total = _page_pixel_count(page, dpi)
    if total > max_page_pixels:
        raise ValueError(
            f"ページ{page_no + 1}: 背景画像が大きすぎます "
            f"(約{total / 1e6:.1f}MP、上限 {max_page_pixels / 1e6:.1f}MP)。"
            "dpiを下げるか max_page_pixels を緩めてください"
        )


def _check_total_pixel_budget(
    doc: pymupdf.Document, dpi: int, max_total_pixels: float
) -> None:
    """全ページ合計の画素数が上限を超えないか、変換開始前に検査する。"""
    total = sum(_page_pixel_count(p, dpi) for p in doc)
    if total > max_total_pixels:
        raise ValueError(
            f"全ページ合計の背景画像が大きすぎます "
            f"(約{total / 1e6:.1f}MP、上限 {max_total_pixels / 1e6:.1f}MP)。"
            "dpiを下げるか、ページ数・max_total_pixelsを見直してください"
        )


def _check_slide_size_limits(width_pt: float, height_pt: float) -> None:
    """1ページ目のサイズが PowerPoint のスライドサイズ制限内かを検査する。

    スライドサイズは1ページ目のサイズで決まるため、これを満たしていれば
    デッキ全体としてPowerPointで開ける寸法になる。
    """
    w_in = width_pt / PT_PER_INCH
    h_in = height_pt / PT_PER_INCH
    if not (PPTX_MIN_SLIDE_INCH <= w_in <= PPTX_MAX_SLIDE_INCH):
        raise ValueError(
            f"PDFページの幅がPowerPointのスライドサイズ制限"
            f"({PPTX_MIN_SLIDE_INCH:.0f}〜{PPTX_MAX_SLIDE_INCH:.0f}インチ)"
            f"を超えています: {w_in:.2f}インチ"
        )
    if not (PPTX_MIN_SLIDE_INCH <= h_in <= PPTX_MAX_SLIDE_INCH):
        raise ValueError(
            f"PDFページの高さがPowerPointのスライドサイズ制限"
            f"({PPTX_MIN_SLIDE_INCH:.0f}〜{PPTX_MAX_SLIDE_INCH:.0f}インチ)"
            f"を超えています: {h_in:.2f}インチ"
        )


def process_page(
    page: pymupdf.Page,
    page_no: int,
    dpi: int,
    max_page_pixels: int,
    ocr: str = DEFAULT_OCR_ENGINE,
    ocr_lang: str = DEFAULT_OCR_LANG,
    ocr_min_conf: float = DEFAULT_OCR_MIN_CONF,
    ocr_timeout: float = DEFAULT_OCR_TIMEOUT,
) -> PageResult:
    result = PageResult(width=page.rect.width, height=page.rect.height)
    _check_pixel_budget(page, dpi, max_page_pixels, page_no)
    ocr = _normalize_ocr_engine(ocr)

    if page.rotation != 0:
        # ページ全体の回転は座標系の補正が絡むため、Phase 1では編集対象に
        # せず安全側に倒す（背景画像はMuPDFが回転込みで正しく描画する）。
        result.warnings.append(
            f"ページの回転({page.rotation}度)を検出しました。"
            "Phase 1では編集対象外とし、背景画像のみ出力します"
        )
        result.image_png = render_background(page, dpi)
        if ocr == OCR_ENGINE_TESSERACT:
            ocr_lines = ocr_image_to_lines(
                result.image_png,
                dpi=dpi,
                lang=ocr_lang,
                min_conf=ocr_min_conf,
                timeout=ocr_timeout,
            )
            result.lines.extend(ocr_lines)
            if ocr_lines:
                result.warnings.append(
                    f"OCRで {len(ocr_lines)} 行の編集可能テキストを追加しました"
                    "（認識誤り・二重写りの可能性があります）"
                )
            else:
                result.warnings.append("OCRを試しましたが文字を検出できませんでした")
        return result

    text_dict = _get_text_dict(page)
    any_visible, any_invisible = _text_visibility(text_dict)
    image_coverage = _image_coverage_ratio(page)
    full_page_image = image_coverage >= FULL_IMAGE_COVERAGE_THRESHOLD

    if not any_visible:
        if any_invisible:
            result.warnings.append(
                "不可視のテキスト層を検出しました（OCR結果等の可能性）。"
                "編集対象にはせず背景画像のみ出力します"
            )
        elif page.get_images():
            result.warnings.append(
                "テキスト層がありません（スキャンPDFの可能性）。背景画像のみ出力します"
            )
        else:
            result.warnings.append("テキストがないページです。背景画像のみ出力します")
    elif full_page_image:
        # 可視テキストはあるが、ページの大部分が画像で覆われている。
        # 画像に焼き込まれた文字を編集可能テキスト化すると、画像と重なって
        # 二重写りになるリスクがあるため、このページ全体を背景画像のみにする。
        result.warnings.append(
            f"ページ面積の約{image_coverage:.0%}が画像で覆われています。"
            "画像上の文字が二重写りになるリスクがあるため、"
            "編集対象にはせず背景画像のみ出力します"
        )
    else:
        result.lines, warns = extract_editable_lines(page, text_dict=text_dict)
        result.warnings.extend(warns)
        redact_text(page, result.lines)
        # redaction漏れの自己検査: 編集対象の文字が背景に残っていれば座標系の不整合
        leftover = page.get_text()
        ghosts = sum(
            1
            for line in result.lines
            for span in line.spans
            if span["text"].strip() and span["text"].strip() in leftover
        )
        if ghosts:
            result.warnings.append(
                f"redactionで消えなかったテキストが {ghosts} 箇所あります（二重写りの可能性）"
            )

    result.image_png = render_background(page, dpi)
    if ocr == OCR_ENGINE_TESSERACT and not result.lines:
        ocr_lines = ocr_image_to_lines(
            result.image_png,
            dpi=dpi,
            lang=ocr_lang,
            min_conf=ocr_min_conf,
            timeout=ocr_timeout,
        )
        result.lines.extend(ocr_lines)
        if ocr_lines:
            result.warnings.append(
                f"OCRで {len(ocr_lines)} 行の編集可能テキストを追加しました"
                "（認識誤り・二重写りの可能性があります）"
            )
        else:
            result.warnings.append("OCRを試しましたが文字を検出できませんでした")
    return result


# ---------------------------------------------------------------------------
# PPTX 生成
# ---------------------------------------------------------------------------

def _pt_to_emu(v: float) -> int:
    return round(v * EMU_PER_PT)


def _set_run_font(run, family: str) -> None:
    """欧文 (a:latin) に加えて日本語 (a:ea) のフォントも設定する。"""
    run.font.name = family  # a:latin
    rPr = run._r.get_or_add_rPr()
    latin = rPr.find(qn("a:latin"))
    anchor = latin
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            anchor.addnext(el)
        el.set("typeface", family)
        anchor = el


def _add_line_textbox(slide, line: Line, scale: float, dx: float, dy: float) -> None:
    """1行分の編集可能テキストボックスをスライドに追加する。"""
    x = line.bbox.x0 * scale + dx
    y = line.bbox.y0 * scale + dy
    w = line.bbox.width * scale + BOX_PAD
    h = line.bbox.height * scale + BOX_PAD

    box = slide.shapes.add_textbox(
        _pt_to_emu(x), _pt_to_emu(y), _pt_to_emu(w), _pt_to_emu(h)
    )
    tf = box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = False
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.TOP

    para = tf.paragraphs[0]
    para.space_before = Pt(0)
    para.space_after = Pt(0)

    for span in line.spans:
        run = para.add_run()
        run.text = span["text"]
        font = run.font
        font.size = Pt(max(span["size"] * scale, 1.0))
        c = span.get("color", 0)
        font.color.rgb = RGBColor((c >> 16) & 0xFF, (c >> 8) & 0xFF, c & 0xFF)
        font.bold = _is_bold(span)
        font.italic = _is_italic(span)
        _set_run_font(run, map_font(span.get("font", ""), span["text"]))


def _add_page_slide(
    prs: Presentation, blank_layout, pr: PageResult, slide_w: float, slide_h: float,
    page_index: int,
) -> list[str]:
    """1ページ分のスライド(背景画像+編集可能テキスト)を追加する。戻り値は警告。"""
    warnings: list[str] = []
    slide = prs.slides.add_slide(blank_layout)

    # ページサイズ混在: 1ページ目基準のスライドに中央配置で収める。
    # scaleは1.0を上限とし、小さいページを勝手に拡大しない。
    if (
        abs(pr.width - slide_w) > PAGE_SIZE_TOL
        or abs(pr.height - slide_h) > PAGE_SIZE_TOL
    ):
        scale = min(1.0, slide_w / pr.width, slide_h / pr.height)
        action = (
            f"{scale:.0%}に縮小して中央配置します"
            if scale < 1.0 - 1e-9
            else "拡大はせず中央配置します"
        )
        warnings.append(
            f"ページ{page_index + 1}: サイズが1ページ目と異なります"
            f"（{pr.width:.0f}x{pr.height:.0f}pt）。{action}"
        )
    else:
        scale = 1.0
    dx = (slide_w - pr.width * scale) / 2
    dy = (slide_h - pr.height * scale) / 2

    slide.shapes.add_picture(
        io.BytesIO(pr.image_png),
        _pt_to_emu(dx),
        _pt_to_emu(dy),
        width=_pt_to_emu(pr.width * scale),
        height=_pt_to_emu(pr.height * scale),
    )
    for line in pr.lines:
        _add_line_textbox(slide, line, scale, dx, dy)

    return warnings


# ---------------------------------------------------------------------------
# 変換パイプライン / CLI
# ---------------------------------------------------------------------------

def _check_positive(name: str, value: float) -> None:
    if value <= 0:
        raise ValueError(f"{name} は正の値である必要があります: {value}")


def looks_like_pdf(path: Path) -> bool:
    """先頭付近に %PDF- マジックバイトがあるかを確認する（簡易な入力検証）。

    PDF仕様上、ヘッダの前に最大1024バイト程度のゴミが許容されるため、
    先頭1024バイトの範囲で探す。app.py（Webアップロード時の検証）からも
    再利用するため公開関数にしている。
    """
    try:
        with open(path, "rb") as f:
            head = f.read(1024)
    except OSError:
        return False
    return b"%PDF-" in head


def check_distinct_input_output(input_pdf: Path, output_pptx: Path) -> None:
    """入力と出力が同一ファイルでないことを検査する（パス一致・ハードリンク双方）。

    worker.py は convert.py を別プロセスで起動する際、実際の出力先ではなく
    一時ディレクトリ内のパスを子プロセスに渡すため、この関数内での検査だけでは
    worker.py 経由の実行を保護できない。worker.py 側でも、子プロセスに渡す前の
    「本来の入出力パス」に対して本関数を呼び出す必要がある。
    """
    if input_pdf.resolve() == output_pptx.resolve():
        raise ValueError(f"入力と出力に同じパスは指定できません: {input_pdf}")
    if output_pptx.exists() and os.path.samefile(input_pdf, output_pptx):
        raise ValueError(
            f"入力と出力が同一ファイルです（ハードリンク等の可能性）: {input_pdf}"
        )


def convert(
    input_pdf: Path,
    output_pptx: Path,
    dpi: int = 150,
    debug_dir: Path | None = None,
    max_pages: int = DEFAULT_MAX_PAGES,
    max_dpi: int = DEFAULT_MAX_DPI,
    max_page_pixels: int = DEFAULT_MAX_PAGE_PIXELS,
    max_total_pixels: float = DEFAULT_MAX_TOTAL_PIXELS,
    max_file_size_mb: float = DEFAULT_MAX_FILE_SIZE_MB,
    max_output_size_mb: float = DEFAULT_MAX_OUTPUT_SIZE_MB,
    timeout_seconds: float | None = None,
    ocr: str = DEFAULT_OCR_ENGINE,
    ocr_lang: str = DEFAULT_OCR_LANG,
    ocr_min_conf: float = DEFAULT_OCR_MIN_CONF,
    ocr_timeout: float = DEFAULT_OCR_TIMEOUT,
) -> list[str]:
    """PDFをPPTXに変換する。戻り値は警告メッセージのリスト。

    制限を超えた場合は処理を始める前、または各ページ処理の開始時点で
    ValueError / TimeoutError を送出して安全に停止する。
    これはプロセス内のソフトタイムアウトであり、ネイティブコード内で
    ハングした場合等に確実に止めたい場合は worker.py の別プロセス+
    ハードタイムアウト(SIGKILL)を使うこと。
    """
    # --- パス・パラメータの検証（重い処理の前に済ませる） ---
    check_distinct_input_output(input_pdf, output_pptx)

    _check_positive("dpi", dpi)
    _check_positive("max_pages", max_pages)
    _check_positive("max_dpi", max_dpi)
    _check_positive("max_page_pixels", max_page_pixels)
    _check_positive("max_total_pixels", max_total_pixels)
    _check_positive("max_file_size_mb", max_file_size_mb)
    _check_positive("max_output_size_mb", max_output_size_mb)
    _check_positive("ocr_timeout", ocr_timeout)
    if timeout_seconds is not None and timeout_seconds < 0:
        raise ValueError(f"timeout_seconds は0以上である必要があります: {timeout_seconds}")
    if not (0 <= ocr_min_conf <= 100):
        raise ValueError(f"ocr_min_conf は0〜100の範囲である必要があります: {ocr_min_conf}")
    ocr = _normalize_ocr_engine(ocr)

    if dpi > max_dpi:
        raise ValueError(f"dpi={dpi} が上限 {max_dpi} を超えています")

    check_ocr_available(ocr, ocr_lang)

    file_size_mb = input_pdf.stat().st_size / (1024 * 1024)
    if file_size_mb > max_file_size_mb:
        raise ValueError(
            f"入力PDFが大きすぎます: {file_size_mb:.1f}MB (上限 {max_file_size_mb}MB)"
        )

    if not looks_like_pdf(input_pdf):
        raise ValueError(
            f"入力ファイルがPDFではないようです（%PDFヘッダが見つかりません）: {input_pdf}"
        )

    if debug_dir is not None and debug_dir.exists() and any(debug_dir.iterdir()):
        raise ValueError(
            f"デバッグ出力先が空ではありません: {debug_dir}"
            "（既存ファイルの上書き事故を防ぐため、空のディレクトリのみ指定できます）"
        )

    warnings: list[str] = []
    start_time = time.monotonic()
    doc = pymupdf.open(input_pdf)
    try:
        if not doc.is_pdf:
            raise ValueError(f"PDFとして認識できませんでした: {input_pdf}")
        if doc.needs_pass:
            raise ValueError(f"パスワード付きPDFは扱えません: {input_pdf}")
        if doc.page_count == 0:
            raise ValueError(f"ページがありません: {input_pdf}")
        if doc.page_count > max_pages:
            raise ValueError(
                f"ページ数が上限を超えています: {doc.page_count} (上限 {max_pages})"
            )

        # 全ページ合計の画素数を、実際に描画する前に検査する。
        _check_total_pixel_budget(doc, dpi, max_total_pixels)

        first_rect = doc[0].rect
        _check_slide_size_limits(first_rect.width, first_rect.height)
        slide_w, slide_h = first_rect.width, first_rect.height

        prs = Presentation()
        prs.slide_width = Emu(_pt_to_emu(slide_w))
        prs.slide_height = Emu(_pt_to_emu(slide_h))
        blank_layout = prs.slide_layouts[6]

        # 各ページの背景画像を溜め込まず、都度スライドに追加してから破棄する
        # （全ページ分をメモリに保持するピーク使用量を抑えるため）。
        for i, page in enumerate(doc):
            if (
                timeout_seconds is not None
                and time.monotonic() - start_time > timeout_seconds
            ):
                raise TimeoutError(
                    f"処理時間が上限({timeout_seconds}秒)を超えました"
                    f"（{i}/{doc.page_count} ページ処理済み）"
                )
            pr = process_page(
                page,
                i,
                dpi,
                max_page_pixels,
                ocr=ocr,
                ocr_lang=ocr_lang,
                ocr_min_conf=ocr_min_conf,
                ocr_timeout=ocr_timeout,
            )
            warnings.extend(f"ページ{i + 1}: {w}" for w in pr.warnings)
            if debug_dir:
                debug_dir.mkdir(parents=True, exist_ok=True)
                (debug_dir / f"page{i + 1:03d}_bg.png").write_bytes(pr.image_png)
            warnings.extend(_add_page_slide(prs, blank_layout, pr, slide_w, slide_h, i))
    finally:
        doc.close()

    prs.save(output_pptx)

    output_size_mb = output_pptx.stat().st_size / (1024 * 1024)
    if output_size_mb > max_output_size_mb:
        output_pptx.unlink(missing_ok=True)
        raise ValueError(
            f"出力PPTXが大きすぎます: {output_size_mb:.1f}MB (上限 {max_output_size_mb}MB)"
        )

    return warnings


def build_arg_parser() -> argparse.ArgumentParser:
    """CLI引数パーサ。worker.py からも共有して二重メンテを避ける。"""
    parser = argparse.ArgumentParser(
        description="PDFを編集可能なPPTXに変換する (Phase 1.1)"
    )
    parser.add_argument("input", type=Path, help="入力PDF")
    parser.add_argument("output", type=Path, help="出力PPTX")
    parser.add_argument("--dpi", type=int, default=150, help="背景画像の解像度 (既定: 150)")
    parser.add_argument(
        "--debug-dir", type=Path, default=None,
        help="redaction後の背景PNGを保存する空のディレクトリ（検証用。"
             "機密PDFや本番では使わないこと）",
    )
    parser.add_argument(
        "--max-pages", type=int, default=DEFAULT_MAX_PAGES,
        help=f"ページ数の上限 (既定: {DEFAULT_MAX_PAGES})",
    )
    parser.add_argument(
        "--max-dpi", type=int, default=DEFAULT_MAX_DPI,
        help=f"dpiの上限 (既定: {DEFAULT_MAX_DPI})",
    )
    parser.add_argument(
        "--max-page-pixels", type=int, default=DEFAULT_MAX_PAGE_PIXELS,
        help=f"背景画像1ページあたりの画素数上限 (既定: {DEFAULT_MAX_PAGE_PIXELS})",
    )
    parser.add_argument(
        "--max-total-pixels", type=float, default=DEFAULT_MAX_TOTAL_PIXELS,
        help=f"背景画像の全ページ合計の画素数上限 (既定: {DEFAULT_MAX_TOTAL_PIXELS})",
    )
    parser.add_argument(
        "--max-file-size-mb", type=float, default=DEFAULT_MAX_FILE_SIZE_MB,
        help=f"入力PDFのファイルサイズ上限MB (既定: {DEFAULT_MAX_FILE_SIZE_MB})",
    )
    parser.add_argument(
        "--max-output-size-mb", type=float, default=DEFAULT_MAX_OUTPUT_SIZE_MB,
        help=f"出力PPTXのファイルサイズ上限MB (既定: {DEFAULT_MAX_OUTPUT_SIZE_MB})",
    )
    parser.add_argument(
        "--timeout", type=float, default=None,
        help="処理時間の上限（秒。ソフトタイムアウト）。既定は無制限",
    )
    parser.add_argument(
        "--ocr",
        choices=OCR_ENGINES,
        default=DEFAULT_OCR_ENGINE,
        help=(
            "OCRエンジン。off=OCRなし、tesseract=ローカルTesseractで"
            f"背景のみページにOCRテキストを重ねる (既定: {DEFAULT_OCR_ENGINE})"
        ),
    )
    parser.add_argument(
        "--ocr-lang",
        default=DEFAULT_OCR_LANG,
        help=f"Tesseract OCRの言語指定 (既定: {DEFAULT_OCR_LANG})",
    )
    parser.add_argument(
        "--ocr-min-conf",
        type=float,
        default=DEFAULT_OCR_MIN_CONF,
        help=f"OCR単語を採用する最低信頼度0〜100 (既定: {DEFAULT_OCR_MIN_CONF})",
    )
    parser.add_argument(
        "--ocr-timeout",
        type=float,
        default=DEFAULT_OCR_TIMEOUT,
        help=f"OCRの1ページあたりタイムアウト秒 (既定: {DEFAULT_OCR_TIMEOUT})",
    )
    return parser


def main(argv: list[str] | None = None) -> int:
    parser = build_arg_parser()
    args = parser.parse_args(argv)

    if not args.input.exists():
        print(f"エラー: 入力ファイルがありません: {args.input}", file=sys.stderr)
        return 1

    try:
        warnings = convert(
            args.input, args.output, dpi=args.dpi, debug_dir=args.debug_dir,
            max_pages=args.max_pages, max_dpi=args.max_dpi,
            max_page_pixels=args.max_page_pixels,
            max_total_pixels=args.max_total_pixels,
            max_file_size_mb=args.max_file_size_mb,
            max_output_size_mb=args.max_output_size_mb,
            timeout_seconds=args.timeout,
            ocr=args.ocr,
            ocr_lang=args.ocr_lang,
            ocr_min_conf=args.ocr_min_conf,
            ocr_timeout=args.ocr_timeout,
        )
    except Exception as e:  # noqa: BLE001 - CLIの最上位でまとめて報告する
        print(f"エラー: {e}", file=sys.stderr)
        return 1

    for w in warnings:
        print(f"警告: {w}", file=sys.stderr)
    print(f"完了: {args.output}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
