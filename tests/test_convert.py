"""convert.py のテスト（合成PDFのみ・ネットワーク不使用）。"""

from __future__ import annotations

import os
import subprocess
import sys
import unicodedata
from pathlib import Path

import pymupdf
import pytest
from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Emu

import convert
import fixtures_gen as fx

ROOT = Path(__file__).resolve().parents[1]
EMU_PER_PT = 12700


# ---------------------------------------------------------------------------
# fixtures
# ---------------------------------------------------------------------------

@pytest.fixture(scope="module")
def business(tmp_path_factory):
    """business.pdf を1回だけ変換し、(pdf, pptx, warnings) を返す。"""
    d = tmp_path_factory.mktemp("business")
    pdf = d / "business.pdf"
    pptx = d / "business.pptx"
    fx.make_business_pdf(pdf)
    warnings = convert.convert(pdf, pptx)
    return pdf, pptx, warnings


def _norm(s: str) -> str:
    """比較用にテキストをNFKC正規化する。

    CJKフォント(特にNoto Sans CJK等のCFF系)ではPyMuPDFのテキスト抽出が
    フォント依存で次のような「見た目は同じだが別コードポイント」の文字を
    返すことがある(macOSのArial Unicodeでは発生しない、Ubuntu CI環境で
    診断して判明した既知挙動):
      - ASCIIスペース(U+0020) -> ノーブレークスペース(U+00A0)
      - 一部の漢字(例: 年 U+5E74) -> CJK互換漢字(例: U+F98E)
    どちらもUnicodeのNFKC正規化で解決できる。これはconvert.py側の不具合
    ではなく元PDFのテキスト層をそのまま反映した結果のため、convert.py側は
    直さずテスト側の比較でのみ正規化する。
    """
    return unicodedata.normalize("NFKC", s)


def _textbox_map(slide):
    """スライド内のテキストボックスを {テキスト: shape} で返す。"""
    return {
        _norm(s.text_frame.text): s
        for s in slide.shapes
        if s.has_text_frame and s.text_frame.text
    }


def _find_box(slide, text):
    boxes = _textbox_map(slide)
    assert text in boxes, f"テキストボックスが見つからない: {text!r} in {list(boxes)}"
    return boxes[text]


# ---------------------------------------------------------------------------
# 基本構造
# ---------------------------------------------------------------------------

def test_slide_size_matches_pdf_page(business):
    pdf, pptx, _ = business
    doc = pymupdf.open(pdf)
    prs = Presentation(pptx)
    assert len(prs.slides._sldIdLst) == doc.page_count == 3
    assert prs.slide_width == round(doc[0].rect.width * EMU_PER_PT)
    assert prs.slide_height == round(doc[0].rect.height * EMU_PER_PT)
    doc.close()


def test_every_slide_has_background_picture(business):
    _, pptx, _ = business
    prs = Presentation(pptx)
    for slide in prs.slides:
        pics = [s for s in slide.shapes if s.shape_type == 13]  # PICTURE
        assert len(pics) == 1
        assert pics[0].left == 0 and pics[0].top == 0


# ---------------------------------------------------------------------------
# テキストの再現
# ---------------------------------------------------------------------------

def test_expected_texts_are_editable(business):
    _, pptx, _ = business
    slide = Presentation(pptx).slides[0]
    boxes = _textbox_map(slide)
    for expected in [fx.TITLE, fx.CONFIDENTIAL, fx.NOTE, *fx.BODY_LINES,
                     "項目", "売上高", "1,234", "営業利益"]:
        assert any(expected in t for t in boxes), f"編集不可: {expected!r}"


def test_position_within_tolerance(business):
    """PPTX上のボックス位置がPDFの行bboxと±3pt以内で一致する。"""
    pdf, pptx, _ = business
    doc = pymupdf.open(pdf)
    lines, _ = convert.extract_editable_lines(doc[0])
    doc.close()
    slide = Presentation(pptx).slides[0]
    boxes = _textbox_map(slide)

    checked = 0
    for line in lines:
        text = _norm("".join(s["text"] for s in line.spans))
        if text not in boxes:
            continue
        shape = boxes[text]
        assert abs(shape.left / EMU_PER_PT - line.bbox.x0) < 3.0, text
        assert abs(shape.top / EMU_PER_PT - line.bbox.y0) < 3.0, text
        checked += 1
    assert checked >= 8  # 主要な行が検査されていること


def test_font_size_color_bold(business):
    _, pptx, _ = business
    slide = Presentation(pptx).slides[0]

    title = _find_box(slide, fx.TITLE).text_frame.paragraphs[0].runs[0]
    assert abs(title.font.size.pt - 20) < 0.5
    assert str(title.font.color.rgb) == "1F3864"
    assert title.font.bold is False

    conf = _find_box(slide, fx.CONFIDENTIAL).text_frame.paragraphs[0].runs[0]
    assert conf.font.bold is True
    assert abs(conf.font.size.pt - 11) < 0.5


def test_textbox_settings(business):
    """内部マージン0・word_wrap無効・上詰めアンカーになっている。"""
    _, pptx, _ = business
    slide = Presentation(pptx).slides[0]
    tf = _find_box(slide, fx.TITLE).text_frame
    assert (tf.margin_left, tf.margin_right, tf.margin_top, tf.margin_bottom) \
        == (Emu(0),) * 4
    assert tf.word_wrap is False
    assert tf.vertical_anchor == MSO_ANCHOR.TOP


def test_ea_font_is_set_for_japanese(business):
    """日本語run に a:ea フォントが設定されている。"""
    _, pptx, _ = business
    from pptx.oxml.ns import qn
    slide = Presentation(pptx).slides[0]
    run = _find_box(slide, fx.TITLE).text_frame.paragraphs[0].runs[0]
    ea = run._r.get_or_add_rPr().find(qn("a:ea"))
    assert ea is not None and ea.get("typeface")


# ---------------------------------------------------------------------------
# 除外対象・検出
# ---------------------------------------------------------------------------

def test_rotated_text_left_in_background(business):
    """回転テキストは編集ボックスにならず、警告が出る。"""
    _, pptx, warnings = business
    slide2 = Presentation(pptx).slides[1]
    boxes = _textbox_map(slide2)
    assert not any(fx.ROTATED_TEXT in t for t in boxes)
    assert any(fx.PAGE2_BODY in t for t in boxes)
    assert any("縦書き・回転テキスト" in w for w in warnings)


def test_scanned_page_detected(business):
    """スキャン風ページは背景画像のみ+警告。"""
    _, pptx, warnings = business
    slide3 = Presentation(pptx).slides[2]
    assert len(slide3.shapes) == 1  # 背景画像のみ
    assert any("スキャン" in w for w in warnings)


def test_no_redaction_ghosts(business):
    """redaction漏れ（二重写り）の自己検査警告が出ていない。"""
    _, _, warnings = business
    assert not any("redaction" in w for w in warnings)


def test_mixed_page_sizes_warns(tmp_path):
    pdf = tmp_path / "mixed.pdf"
    pptx = tmp_path / "mixed.pptx"
    fx.make_mixed_pdf(pdf)
    warnings = convert.convert(pdf, pptx)
    assert any("サイズが1ページ目と異なります" in w for w in warnings)
    prs = Presentation(pptx)
    assert len(prs.slides._sldIdLst) == 2
    # 2ページ目のテキストも（縮小されつつ）編集可能であること
    assert any(fx.MIXED_P2_TEXT in t for t in _textbox_map(prs.slides[1]))


def test_rotated_page_falls_back_to_background(tmp_path):
    """/Rotate 90 付きページはPhase 1では編集対象にせず背景画像のみにする。"""
    pdf = tmp_path / "rot.pdf"
    pptx = tmp_path / "rot.pptx"
    fx.make_rotated_page_pdf(pdf)
    warnings = convert.convert(pdf, pptx)
    slide = Presentation(pptx).slides[0]
    assert not any(fx.ROTATED_PAGE_TEXT in t for t in _textbox_map(slide))
    assert len(slide.shapes) == 1  # 背景画像のみ（テキストボックスなし）
    assert any("回転" in w for w in warnings)
    assert not any("redaction" in w for w in warnings)


def test_small_page_not_enlarged(tmp_path):
    """ページ1より小さいページは拡大せず、実寸のまま中央配置する。"""
    pdf = tmp_path / "small_mixed.pdf"
    pptx = tmp_path / "small_mixed.pptx"
    fx.make_small_mixed_pdf(pdf)
    warnings = convert.convert(pdf, pptx)
    assert any("拡大" in w for w in warnings)

    slide2 = Presentation(pptx).slides[1]
    pics = [s for s in slide2.shapes if s.shape_type == 13]
    assert len(pics) == 1
    assert pics[0].width == convert._pt_to_emu(200)
    assert pics[0].height == convert._pt_to_emu(100)
    # 中央配置なので原点(0,0)には置かれない
    assert pics[0].left > 0 and pics[0].top > 0
    assert any(fx.SMALL_PAGE_TEXT in t for t in _textbox_map(slide2))


# ---------------------------------------------------------------------------
# 不可視OCRテキスト
# ---------------------------------------------------------------------------

def test_invisible_ocr_text_not_editable(tmp_path):
    """不可視OCRテキスト層のみのページは背景画像のみ+警告になる（可視の二重化防止）。"""
    pdf = tmp_path / "ocr.pdf"
    pptx = tmp_path / "ocr.pptx"
    fx.make_invisible_ocr_pdf(pdf)
    warnings = convert.convert(pdf, pptx)

    slide = Presentation(pptx).slides[0]
    assert len(slide.shapes) == 1  # 背景画像のみ
    assert not any(fx.OCR_INVISIBLE_TEXT in t for t in _textbox_map(slide))
    assert any("不可視" in w for w in warnings)


def test_is_invisible_detects_alpha_zero_and_render_mode3():
    doc = pymupdf.open()
    page = doc.new_page(width=200, height=200)
    page.insert_text((10, 50), "visible", fontname="helv", fontsize=12)
    page.insert_text((10, 80), "transparent", fontname="helv", fontsize=12, fill_opacity=0)
    page.insert_text((10, 110), "rendermode3", fontname="helv", fontsize=12, render_mode=3)
    d = convert._get_text_dict(page)
    spans = [
        s
        for b in d["blocks"] if b["type"] == 0
        for l in b["lines"]
        for s in l["spans"]
    ]
    doc.close()
    by_text = {s["text"].strip(): s for s in spans}
    assert convert._is_invisible(by_text["visible"]) is False
    assert convert._is_invisible(by_text["transparent"]) is True
    assert convert._is_invisible(by_text["rendermode3"]) is True


# ---------------------------------------------------------------------------
# フォールバック文字との重なり
# ---------------------------------------------------------------------------

def test_overlapping_horizontal_text_left_in_background(tmp_path):
    """回転テキストとbboxが重なる横書き行はredactionせず背景に残す。

    重ならない通常行は従来通り編集可能になる（過剰な除外になっていないこと）。
    """
    pdf = tmp_path / "overlap.pdf"
    pptx = tmp_path / "overlap.pptx"
    fx.make_overlap_pdf(pdf)
    warnings = convert.convert(pdf, pptx)

    slide = Presentation(pptx).slides[0]
    boxes = _textbox_map(slide)
    assert not any(fx.OVERLAP_HORIZONTAL_TEXT in t for t in boxes), list(boxes)
    assert not any(fx.OVERLAP_ROTATED_TEXT in t for t in boxes), list(boxes)
    # フォント(特にCFF系CJKフォント)によっては行が複数spanに分割されて
    # 抽出されることがあるため、単一ボックスへの部分一致ではなく
    # 全ボックスを連結したテキストに対して判定する（診断用に一覧も表示）。
    all_text = "".join(boxes)
    assert fx.OVERLAP_NORMAL_TEXT in all_text, list(boxes)
    assert any("重なる" in w for w in warnings)
    assert not any("redaction" in w for w in warnings)


# ---------------------------------------------------------------------------
# 入出力パス防御
# ---------------------------------------------------------------------------

def test_convert_rejects_same_input_output_path(tmp_path):
    pdf = tmp_path / "same.pdf"
    fx.make_business_pdf(pdf)
    with pytest.raises(ValueError, match="同じパス"):
        convert.convert(pdf, pdf)
    # 入力ファイルが破壊されていないこと
    assert pymupdf.open(pdf).page_count == 3


def test_cli_rejects_same_input_output_path(tmp_path):
    pdf = tmp_path / "same.pdf"
    fx.make_business_pdf(pdf)
    res = subprocess.run(
        [sys.executable, str(ROOT / "convert.py"), str(pdf), str(pdf)],
        capture_output=True, text=True,
    )
    assert res.returncode == 1
    assert pymupdf.open(pdf).page_count == 3


# ---------------------------------------------------------------------------
# Web公開前提の基本制限
# ---------------------------------------------------------------------------

def test_max_pages_exceeded(tmp_path):
    pdf = tmp_path / "in.pdf"
    pptx = tmp_path / "out.pptx"
    fx.make_business_pdf(pdf)  # 3ページ
    with pytest.raises(ValueError, match="ページ数"):
        convert.convert(pdf, pptx, max_pages=2)


def test_max_dpi_exceeded(tmp_path):
    pdf = tmp_path / "in.pdf"
    pptx = tmp_path / "out.pptx"
    fx.make_business_pdf(pdf)
    with pytest.raises(ValueError, match="dpi"):
        convert.convert(pdf, pptx, dpi=500, max_dpi=300)


def test_max_file_size_exceeded(tmp_path):
    pdf = tmp_path / "in.pdf"
    pptx = tmp_path / "out.pptx"
    fx.make_business_pdf(pdf)
    with pytest.raises(ValueError, match="大きすぎます"):
        convert.convert(pdf, pptx, max_file_size_mb=0.0001)


def test_max_page_pixels_exceeded(tmp_path):
    pdf = tmp_path / "in.pdf"
    pptx = tmp_path / "out.pptx"
    fx.make_business_pdf(pdf)
    with pytest.raises(ValueError, match="大きすぎます"):
        convert.convert(pdf, pptx, max_page_pixels=100)


def test_timeout_exceeded(tmp_path):
    pdf = tmp_path / "in.pdf"
    pptx = tmp_path / "out.pptx"
    fx.make_business_pdf(pdf)
    with pytest.raises(TimeoutError):
        convert.convert(pdf, pptx, timeout_seconds=0.0)


def test_within_limits_still_succeeds(tmp_path):
    """既定の制限は通常のビジネス文書PDFを妨げない。"""
    pdf = tmp_path / "in.pdf"
    pptx = tmp_path / "out.pptx"
    fx.make_business_pdf(pdf)
    warnings = convert.convert(pdf, pptx)
    assert pptx.exists()
    assert not any("大きすぎます" in w or "上限" in w for w in warnings)


# ---------------------------------------------------------------------------
# 画像が多いPDFでのget_text負荷
# ---------------------------------------------------------------------------

def test_image_heavy_pdf_text_dict_excludes_images(tmp_path):
    """TEXT_PRESERVE_IMAGESを外しているため、画像ブロックが混入しない。"""
    pdf = tmp_path / "images.pdf"
    fx.make_image_heavy_pdf(pdf, n_images=15)
    doc = pymupdf.open(pdf)
    page = doc[0]

    d = convert._get_text_dict(page)

    # 実行時間のassertは共有CIランナーで揺れるため行わず、画像ブロックが
    # 実際に混入していないこと（TEXT_PRESERVE_IMAGESを外した効果）だけを
    # 確定的に検証する。
    assert not any(b["type"] == 1 for b in d["blocks"])  # 画像ブロックが無い

    lines, _ = convert.extract_editable_lines(page, text_dict=d)
    assert any(fx.PAGE2_BODY in "".join(s["text"] for s in line.spans) for line in lines)
    doc.close()


def test_image_heavy_pdf_converts_correctly(tmp_path):
    pdf = tmp_path / "images.pdf"
    pptx = tmp_path / "images.pptx"
    fx.make_image_heavy_pdf(pdf, n_images=15)
    warnings = convert.convert(pdf, pptx)
    slide = Presentation(pptx).slides[0]
    assert any(fx.PAGE2_BODY in t for t in _textbox_map(slide))
    assert not any("redaction" in w for w in warnings)


# ---------------------------------------------------------------------------
# hard link対策
# ---------------------------------------------------------------------------

def test_hardlink_input_output_rejected(tmp_path):
    """resolve()では区別できないハードリンクも samefile で拒否する。"""
    pdf = tmp_path / "in.pdf"
    linked = tmp_path / "out.pptx"
    fx.make_business_pdf(pdf)
    os.link(pdf, linked)  # 同一inodeへの別名（resolve()の結果はそれぞれ異なる）
    assert pdf.resolve() != linked.resolve()

    with pytest.raises(ValueError, match="同一ファイル"):
        convert.convert(pdf, linked)
    # 入力ファイルが破壊されていないこと
    assert pymupdf.open(pdf).page_count == 3


# ---------------------------------------------------------------------------
# debug-dir安全化
# ---------------------------------------------------------------------------

def test_debug_dir_with_existing_files_rejected(tmp_path):
    pdf = tmp_path / "in.pdf"
    pptx = tmp_path / "out.pptx"
    debug_dir = tmp_path / "debug"
    debug_dir.mkdir()
    (debug_dir / "stale.png").write_bytes(b"leftover")
    fx.make_business_pdf(pdf)

    with pytest.raises(ValueError, match="空ではありません"):
        convert.convert(pdf, pptx, debug_dir=debug_dir)
    assert not pptx.exists()
    assert (debug_dir / "stale.png").exists()  # 既存ファイルには触れない


def test_debug_dir_empty_or_new_is_accepted(tmp_path):
    pdf = tmp_path / "in.pdf"
    pptx = tmp_path / "out.pptx"
    debug_dir = tmp_path / "debug"  # 未作成
    fx.make_business_pdf(pdf)
    convert.convert(pdf, pptx, debug_dir=debug_dir)
    assert (debug_dir / "page001_bg.png").exists()


# ---------------------------------------------------------------------------
# 総リソース上限
# ---------------------------------------------------------------------------

def test_max_total_pixels_exceeded(tmp_path):
    """1ページごとは上限内でも、全ページ合計が上限を超えれば変換前に停止する。"""
    pdf = tmp_path / "in.pdf"
    pptx = tmp_path / "out.pptx"
    fx.make_business_pdf(pdf)  # 3ページ、1ページはデフォルト上限内
    with pytest.raises(ValueError, match="合計"):
        convert.convert(pdf, pptx, max_total_pixels=100)
    assert not pptx.exists()


def test_max_output_size_exceeded_removes_file(tmp_path):
    pdf = tmp_path / "in.pdf"
    pptx = tmp_path / "out.pptx"
    fx.make_business_pdf(pdf)
    with pytest.raises(ValueError, match="出力PPTX"):
        convert.convert(pdf, pptx, max_output_size_mb=0.0001)
    assert not pptx.exists()  # 超過分は削除される


# ---------------------------------------------------------------------------
# 全面画像ページの安全フォールバック
# ---------------------------------------------------------------------------

def test_full_page_image_with_visible_text_falls_back(tmp_path):
    """全面画像 + 可視テキストは二重写りリスクのため編集対象にしない。"""
    pdf = tmp_path / "full_image.pdf"
    pptx = tmp_path / "full_image.pptx"
    fx.make_full_image_visible_text_pdf(pdf)
    warnings = convert.convert(pdf, pptx)

    slide = Presentation(pptx).slides[0]
    assert len(slide.shapes) == 1  # 背景画像のみ
    assert not any(fx.FULL_IMAGE_VISIBLE_TEXT in t for t in _textbox_map(slide))
    assert any("画像で覆われています" in w for w in warnings)


def test_image_coverage_ratio_detects_full_page_image():
    doc = pymupdf.open()
    page = doc.new_page(width=200, height=200)
    from PIL import Image as PILImage
    import io as _io
    img = PILImage.new("RGB", (50, 50), (1, 2, 3))
    buf = _io.BytesIO()
    img.save(buf, format="PNG")
    page.insert_image(page.rect, stream=buf.getvalue())
    ratio = convert._image_coverage_ratio(page)
    doc.close()
    assert ratio > 0.99


# ---------------------------------------------------------------------------
# OCR（任意機能）
# ---------------------------------------------------------------------------

def test_ocr_tsv_coordinates_are_converted_to_pdf_points():
    tsv = (
        "level\tpage_num\tblock_num\tpar_num\tline_num\tword_num\t"
        "left\ttop\twidth\theight\tconf\ttext\n"
        "5\t1\t1\t1\t1\t1\t144\t72\t72\t24\t92\tHello\n"
        "5\t1\t1\t1\t1\t2\t220\t72\t80\t24\t90\tWorld\n"
    )
    lines = convert._ocr_lines_from_tsv(tsv, dpi=144, min_conf=35)
    assert len(lines) == 1
    assert lines[0].spans[0]["text"] == "Hello World"
    assert round(lines[0].bbox.x0, 2) == 72.0
    assert round(lines[0].bbox.y0, 2) == 36.0


def test_ocr_adds_editable_text_to_scan_page(tmp_path, monkeypatch):
    pdf = tmp_path / "scan.pdf"
    pptx = tmp_path / "scan.pptx"
    fx.make_business_pdf(pdf)

    def _available(*args, **kwargs):
        return None

    def _fake_ocr(*args, **kwargs):
        bbox = pymupdf.Rect(60, 120, 220, 140)
        return [
            convert.Line(
                bbox=bbox,
                spans=[
                    {
                        "text": "OCR TEST TEXT",
                        "bbox": tuple(bbox),
                        "font": convert.FALLBACK_LATIN_SANS,
                        "size": 12,
                        "color": 0,
                        "flags": 0,
                    }
                ],
            )
        ]

    monkeypatch.setattr(convert, "check_ocr_available", _available)
    monkeypatch.setattr(convert, "ocr_image_to_lines", _fake_ocr)

    warnings = convert.convert(pdf, pptx, ocr=convert.OCR_ENGINE_TESSERACT)
    slide3 = Presentation(pptx).slides[2]
    assert "OCR TEST TEXT" in _textbox_map(slide3)
    assert any("OCRで" in w for w in warnings)


def test_ocr_missing_tesseract_is_clear(monkeypatch):
    monkeypatch.setattr(convert.shutil, "which", lambda name: None)
    with pytest.raises(ValueError, match="Tesseract"):
        convert.check_ocr_available(convert.OCR_ENGINE_TESSERACT)


# ---------------------------------------------------------------------------
# 入力検証
# ---------------------------------------------------------------------------

def test_non_pdf_file_rejected(tmp_path):
    fake = tmp_path / "not_a.pdf"
    fake.write_bytes(b"this is not a pdf file, just plain text padding" * 10)
    pptx = tmp_path / "out.pptx"
    with pytest.raises(ValueError, match="PDFではないよう"):
        convert.convert(fake, pptx)


@pytest.mark.parametrize(
    "kwargs",
    [
        {"dpi": 0},
        {"dpi": -10},
        {"max_pages": 0},
        {"max_dpi": -1},
        {"max_page_pixels": 0},
        {"max_total_pixels": -1},
        {"max_file_size_mb": 0},
        {"max_output_size_mb": -5},
    ],
)
def test_non_positive_limits_rejected(tmp_path, kwargs):
    pdf = tmp_path / "in.pdf"
    pptx = tmp_path / "out.pptx"
    fx.make_business_pdf(pdf)
    with pytest.raises(ValueError, match="正の値"):
        convert.convert(pdf, pptx, **kwargs)


def test_negative_timeout_rejected(tmp_path):
    pdf = tmp_path / "in.pdf"
    pptx = tmp_path / "out.pptx"
    fx.make_business_pdf(pdf)
    with pytest.raises(ValueError, match="timeout_seconds"):
        convert.convert(pdf, pptx, timeout_seconds=-1.0)


def test_extreme_large_page_size_rejected(tmp_path):
    pdf = tmp_path / "huge.pdf"
    pptx = tmp_path / "out.pptx"
    fx.make_extreme_page_size_pdf(pdf)
    with pytest.raises(ValueError, match="スライドサイズ制限"):
        convert.convert(pdf, pptx)


def test_extreme_tiny_page_size_rejected(tmp_path):
    pdf = tmp_path / "tiny.pdf"
    pptx = tmp_path / "out.pptx"
    fx.make_tiny_page_size_pdf(pdf)
    with pytest.raises(ValueError, match="スライドサイズ制限"):
        convert.convert(pdf, pptx)


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def test_cli_smoke(tmp_path):
    pdf = tmp_path / "in.pdf"
    out = tmp_path / "out.pptx"
    fx.make_business_pdf(pdf)
    res = subprocess.run(
        [sys.executable, str(ROOT / "convert.py"), str(pdf), str(out),
         "--debug-dir", str(tmp_path / "debug")],
        capture_output=True, text=True,
    )
    assert res.returncode == 0, res.stderr
    assert out.exists()
    assert (tmp_path / "debug" / "page001_bg.png").exists()
    assert "完了" in res.stdout


def test_cli_missing_input(tmp_path):
    res = subprocess.run(
        [sys.executable, str(ROOT / "convert.py"),
         str(tmp_path / "nai.pdf"), str(tmp_path / "out.pptx")],
        capture_output=True, text=True,
    )
    assert res.returncode == 1
